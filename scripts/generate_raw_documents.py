"""
scripts/generate_raw_documents.py

Comprehensive Document Synthesizer for PragyanAI College Intelligence Hub.
Generates all raw unstructured files in:
  - data/raw/brochures/       (Admissions, Management Quota, Scholarships & ROI PDFs)
  - data/raw/presentations/   (CoE Facilities, R&D Labs & Dept Overview PPTX decks)
  - data/raw/regulatory/      (NAAC SSR Criterion 1-7 & NBA SAR OBE Compliance PDFs)
"""

import os
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==============================================================================
# 0. DIRECTORY SETUP
# ==============================================================================
BASE_DIR = "data"
RAW_BROCHURES = os.path.join(BASE_DIR, "raw", "brochures")
RAW_PRESENTATIONS = os.path.join(BASE_DIR, "raw", "presentations")
RAW_REGULATORY = os.path.join(BASE_DIR, "raw", "regulatory")

for directory in [RAW_BROCHURES, RAW_PRESENTATIONS, RAW_REGULATORY]:
    os.makedirs(directory, exist_ok=True)

print(f"[*] Verified raw document directories under '{BASE_DIR}/raw/'")


# ==============================================================================
# 1. ENTERPRISE PDF ENGINE (FPDF2)
# ==============================================================================
class EnterpriseInstitutionalPDF(FPDF):
    """Custom PDF generator with standardized header, footer, callout boxes, and tables."""

    def __init__(self, doc_category="Official Institutional Document", *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.doc_category = doc_category

    def header(self):
        # Top banner background
        self.set_fill_color(15, 23, 42)  # Dark Slate (Tailwind 900)
        self.rect(0, 0, 210, 18, "F")

        # Header Title & Subtitle
        self.set_xy(10, 4)
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(255, 255, 255)
        self.cell(100, 5, "PRAGYANAI INSTITUTIONAL INTELLIGENCE & ACADEMIC AUDIT", 0, 0, "L")

        self.set_xy(110, 4)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(203, 213, 225)
        self.cell(90, 5, self.doc_category.upper(), 0, 0, "R")

        self.set_xy(10, 10)
        self.set_font("Helvetica", "", 7)
        self.set_text_color(148, 163, 184)
        self.cell(190, 4, "Verified Institutional Document Repository | 2026-27 Academic Cycle", 0, 0, "L")

        self.set_draw_color(37, 99, 235)  # Royal Blue Accent line
        self.set_line_width(0.8)
        self.line(0, 18, 210, 18)
        self.ln(12)

    def footer(self):
        self.set_y(-16)
        self.set_draw_color(226, 232, 240)
        self.set_line_width(0.4)
        self.line(10, self.get_y(), 200, self.get_y())

        self.set_font("Helvetica", "I", 8)
        self.set_text_color(148, 163, 184)
        self.cell(95, 10, "Confidential - Institutional Decision Support Engine", 0, 0, "L")
        self.cell(95, 10, f"Page {self.page_no()}/{{nb}}", 0, 0, "R")

    def add_document_title(self, title: str, subtitle: str = ""):
        self.set_font("Helvetica", "B", 18)
        self.set_text_color(15, 23, 42)
        self.multi_cell(0, 8, title, 0, "L")
        if subtitle:
            self.set_font("Helvetica", "I", 10)
            self.set_text_color(100, 116, 139)
            self.multi_cell(0, 6, subtitle, 0, "L")
        self.ln(4)

    def add_section(self, heading: str, body_text: str):
        self.set_font("Helvetica", "B", 12)
        self.set_text_color(29, 78, 216)  # Deep Blue
        self.cell(0, 7, heading, 0, 1, "L")
        self.set_draw_color(219, 234, 254)
        self.set_line_width(0.3)
        self.line(self.get_x(), self.get_y(), self.get_x() + 190, self.get_y())
        self.ln(2)

        self.set_font("Helvetica", "", 9.5)
        self.set_text_color(51, 65, 85)
        self.multi_cell(0, 5.2, body_text, 0, "L")
        self.ln(3)

    def add_callout(self, title: str, text: str):
        self.set_fill_color(248, 250, 252)
        self.set_draw_color(37, 99, 235)
        self.set_line_width(0.6)
        
        start_y = self.get_y()
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(30, 58, 138)
        
        # Calculate callout box height dynamically
        self.rect(10, start_y, 190, 22, "DF")
        self.set_xy(14, start_y + 2)
        self.cell(180, 5, f"[KEY TAKEAWAY] {title}", 0, 1, "L")
        self.set_xy(14, start_y + 8)
        self.set_font("Helvetica", "", 8.5)
        self.set_text_color(71, 85, 105)
        self.multi_cell(180, 4.5, text, 0, "L")
        self.set_y(start_y + 26)

    def add_table(self, headers: list, rows: list, col_widths: list):
        self.set_font("Helvetica", "B", 8.5)
        self.set_fill_color(241, 245, 249)
        self.set_text_color(15, 23, 42)
        self.set_draw_color(203, 213, 225)
        self.set_line_width(0.2)

        # Header Row
        for idx, header in enumerate(headers):
            self.cell(col_widths[idx], 7, header, 1, 0, "C", fill=True)
        self.ln(7)

        # Data Rows
        self.set_font("Helvetica", "", 8.5)
        self.set_text_color(51, 65, 85)
        for row in rows:
            for idx, item in enumerate(row):
                align = "C" if idx > 0 else "L"
                self.cell(col_widths[idx], 6.5, str(item), 1, 0, align)
            self.ln(6.5)
        self.ln(4)


# ==============================================================================
# 2. GENERATE PDF BROCHURES (data/raw/brochures/)
# ==============================================================================
print("[*] Generating high-fidelity Admission & ROI PDF brochures...")

# A. Management Admission & Fee Structure Guide
pdf_admissions = EnterpriseInstitutionalPDF(doc_category="Admissions & Fee Intelligence")
pdf_admissions.alias_nb_pages()
pdf_admissions.add_page()
pdf_admissions.add_document_title(
    "Management Admissions, Fee Structure & Quota Regulations (2026-27)",
    "Comprehensive Admission Matrix across 15 Benchmark Autonomous Engineering Colleges in Karnataka"
)
pdf_admissions.add_section(
    "1. Overview of Admission Quotas & Regulatory Framework",
    "Under the Karnataka Educational Institutions (Prohibition of Capitation Fee) Act and KEA consensual "
    "agreements, engineering seat distribution in autonomous institutions is divided into three distinct tracks: "
    "Government CET Quota (45%), COMEDK Merit Quota (30%), and Institutional/Management Quota (25%). Management "
    "admissions are strictly governed by transparent eligibility criteria requiring a minimum of 60% aggregate in "
    "Physics, Mathematics, and Chemistry/Computer Science in 10+2 / Pre-University examinations."
)

pdf_admissions.add_table(
    headers=["Branch Name", "CET Govt Fee", "COMEDK Fee", "Mgmt Fee (Tier-1)", "Mgmt Fee (Tier-2)"],
    rows=[
        ["Computer Science & Engineering (CSE)", "INR 1.07 L/yr", "INR 2.81 L/yr", "INR 14.0 - 16.0 L/yr", "INR 7.5 - 9.0 L/yr"],
        ["AI & Data Science (AI-DS)", "INR 1.07 L/yr", "INR 2.81 L/yr", "INR 12.0 - 14.0 L/yr", "INR 6.5 - 8.0 L/yr"],
        ["Information Science & Engg (ISE)", "INR 1.07 L/yr", "INR 2.81 L/yr", "INR 11.0 - 13.0 L/yr", "INR 6.0 - 7.5 L/yr"],
        ["Electronics & Communication (ECE)", "INR 1.07 L/yr", "INR 2.81 L/yr", "INR 9.0 - 11.0 L/yr", "INR 5.0 - 6.5 L/yr"],
        ["Mechanical / Core Automation", "INR 1.07 L/yr", "INR 2.81 L/yr", "INR 4.5 - 6.0 L/yr", "INR 3.0 - 4.5 L/yr"],
    ],
    col_widths=[60, 30, 30, 35, 35]
)

pdf_admissions.add_section(
    "2. Fee Component Bifurcation & Transparent Cost of Attendance",
    "The institutional annual fee includes: (a) Tuition Fee covering academic pedagogy and faculty credits; "
    "(b) University Registration & VTU/Autonomous Exam Fees; (c) Center of Excellence (CoE) Advanced Lab Access; "
    "(d) Value-Added Pre-Placement Bootcamps (Full Stack, GenAI, RTL/VLSI); and (e) Student Medical Insurance and "
    "Digital Library Subscriptions. Hostel accommodation with high-speed Wi-Fi and mess facilities ranges from "
    "INR 1.25 Lakhs to 1.85 Lakhs per academic annum."
)

pdf_admissions.add_callout(
    "Direct Counseling & Lateral Transfer Assistance",
    "Institutional management seats feature direct seat lock-in prior to KEA mop-up rounds. Lateral entry diploma "
    "candidates with 65%+ aggregate are eligible for 2nd-year direct branch admissions under institutional quota."
)

pdf_admissions.add_section(
    "3. Merit Scholarships and Fee Concession Matrix",
    "Institutions offer progressive fee waivers for meritorious candidates: (1) KCET rank < 2,000 or COMEDK rank < 1,500 "
    "qualifies for a 50% waiver on management tuition; (2) National level sports medalists / Olympiad finalists receive a "
    "25% concession; (3) Single girl child and defense ward quotas provide INR 50,000 annual fee subsidies."
)

pdf_admissions.output(os.path.join(RAW_BROCHURES, "Admission_Flyer_2026.pdf"))


# B. Placement ROI & Career Transformation Brochure
pdf_roi = EnterpriseInstitutionalPDF(doc_category="Placement ROI & Career Pathways")
pdf_roi.alias_nb_pages()
pdf_roi.add_page()
pdf_roi.add_document_title(
    "Institutional Placement ROI & Career Transformation Report",
    "Quantitative Salary Benchmarks, Recruiter Density & 4-Year Return on Capital Analysis"
)
pdf_roi.add_section(
    "1. 4-Year Cost of Study vs. First-Year Career Compensation (ROI Index)",
    "The Return on Investment (ROI) index measures total tuition and living capital deployed over 4 years against "
    "median first-year compensation. In benchmark institutions (e.g., RVCE, BMSCE, MSRIT), candidates graduating from "
    "Computing and Electronics branches recover their entire educational capital outlay within 14 to 22 months of "
    "joining product engineering teams."
)

pdf_roi.add_table(
    headers=["Institution", "Median CTC", "Top 10% CTC", "SuperDream (>25 LPA)", "4-Yr ROI Payback"],
    rows=[
        ["RV College of Engineering", "14.50 LPA", "28.50 LPA", "184 Offers", "14 Months"],
        ["BMS College of Engineering", "11.20 LPA", "22.40 LPA", "132 Offers", "18 Months"],
        ["Ramaiah Institute of Tech (MSRIT)", "10.50 LPA", "20.80 LPA", "118 Offers", "19 Months"],
        ["PES University (RR Campus)", "13.00 LPA", "26.00 LPA", "165 Offers", "16 Months"],
        ["Dayananda Sagar CE (DSCE)", "8.20 LPA", "16.50 LPA", "78 Offers", "24 Months"],
        ["JSS Science & Tech Univ (SJCE)", "9.20 LPA", "18.00 LPA", "92 Offers", "21 Months"],
    ],
    col_widths=[55, 30, 30, 45, 30]
)

pdf_roi.add_section(
    "2. Recruiter Tier Breakdown & Hiring Ecosystem",
    "Placement opportunities are classified into 4 distinct compensation categories: Tier-1 SuperDream (INR 25-62 LPA: "
    "Microsoft, Google, Amazon, Apple, Goldman Sachs); Tier-2 Dream Core (INR 12-24 LPA: Qualcomm, Cisco, Samsung R&D, "
    "Texas Instruments, Intel); Tier-3 Product/Fintech (INR 7-12 LPA: Bosch, Dell, Oracle, SAP Labs); and Tier-4 "
    "Enterprise IT (INR 4.5-7 LPA: TCS Digital, Accenture, Infosys, Cognizant)."
)

pdf_roi.add_callout(
    "Pre-Placement Skill Acceleration Engine",
    "Over 88% of SuperDream offer holders completed specialized multi-agent AI, distributed systems, or RTL-VLSI "
    "bootcamps starting from Semester 5, logging 250+ GitHub commits and winning tier-1 hackathons."
)

pdf_roi.output(os.path.join(RAW_BROCHURES, "Placement_ROI_Report_2026.pdf"))
print("  [+] Generated 'Admission_Flyer_2026.pdf' and 'Placement_ROI_Report_2026.pdf'")


# ==============================================================================
# 3. GENERATE REGULATORY & ACCREDITATION DOSSIERS (data/raw/regulatory/)
# ==============================================================================
print("[*] Generating NAAC SSR & NBA OBE Regulatory Compliance Dossiers...")

# A. NAAC Self-Study Report (SSR) Summary Dossier
pdf_naac = EnterpriseInstitutionalPDF(doc_category="NAAC Accreditation Self-Study Dossier")
pdf_naac.alias_nb_pages()
pdf_naac.add_page()
pdf_naac.add_document_title(
    "NAAC Institutional Self-Study Report (SSR) - Executive Audit",
    "Criterion-Wise Quantitative Breakdown for Autonomous Engineering Colleges (Cycle-4)"
)
pdf_naac.add_section(
    "Criterion 1: Curricular Aspects (Institutional Weightage: 150)",
    "1.1 Curriculum Design & Agility: 100% autonomous programs operate on Choice-Based Credit System (CBCS) with Outcome-"
    "Based Education (OBE) frameworks. Curricula are revised triennially with 35% representation from industry leaders.\n"
    "1.2 Value-Added Electives: 68 industry-certified modular electives offered across AI/ML, Quantum Computing, Electric "
    "Vehicle powertrains, and Smart Urban Infrastructure. 94.2% student participation recorded in credit-bearing internships."
)

pdf_naac.add_section(
    "Criterion 2: Teaching-Learning and Evaluation (Weightage: 200)",
    "2.1 Faculty Cadre & Ph.D. Density: 68.4% of full-time faculty hold doctoral degrees from premier institutes (IISc, "
    "IITs, NITs). Student to Full-Time Faculty ratio is maintained at 13.8:1.\n"
    "2.2 Experiential Pedagogy: 100% of laboratory courses utilize project-based learning. Digital smart classrooms "
    "integrated with automated lecture recording and learning management systems (LMS)."
)

pdf_naac.add_table(
    headers=["NAAC Criterion", "Metric Description", "Benchmark Score", "Institutional Attainment"],
    rows=[
        ["Criterion 1", "Curricular Aspects & Academic Flexibility", "150 / 150", "144 / 150 (96.0%)"],
        ["Criterion 2", "Teaching, Learning & Continuous Evaluation", "200 / 200", "191 / 200 (95.5%)"],
        ["Criterion 3", "Research, Innovations & Industrial Consulting", "250 / 250", "236 / 250 (94.4%)"],
        ["Criterion 4", "Infrastructure & Digital Library Resources", "100 / 100", "98 / 100 (98.0%)"],
        ["Criterion 5", "Student Support, Mentorship & Progression", "100 / 100", "94 / 100 (94.0%)"],
        ["Criterion 6", "Governance, Leadership & Financial Strategy", "100 / 100", "95 / 100 (95.0%)"],
        ["Criterion 7", "Institutional Values & Environmental Best Practices", "100 / 100", "97 / 100 (97.0%)"],
    ],
    col_widths=[30, 85, 35, 40]
)

pdf_naac.add_section(
    "Criterion 3: Research, Innovations and Extension (Weightage: 250)",
    "Sponsored R&D projects mobilized INR 24.8 Crores across DST, SERB, DRDO, ISRO, and AICTE grants over the last 3 years. "
    "Institutional innovation incubators have graduated 28 deep-tech student startups, securing over INR 6.5 Crores in "
    "seed capital. Faculty and students filed 58 Indian and international patents, with 19 commercialized."
)

pdf_naac.output(os.path.join(RAW_REGULATORY, "NAAC_Self_Study_Summary.pdf"))


# B. NBA Outcome-Based Education (OBE) Compliance Report
pdf_nba = EnterpriseInstitutionalPDF(doc_category="NBA Tier-1 Accreditation Compliance")
pdf_nba.alias_nb_pages()
pdf_nba.add_page()
pdf_nba.add_document_title(
    "National Board of Accreditation (NBA) - Tier-1 SAR Compliance",
    "Outcome-Based Education (OBE), Program Outcomes (PO1-PO12) & Faculty Cadre Attainment"
)
pdf_nba.add_section(
    "1. Program Educational Objectives (PEOs) & Program Outcomes (POs)",
    "The Department of Computer Science & Engineering and allied branches strictly comply with Washington Accord Tier-1 "
    "standards. Course Outcomes (COs) are formally mapped to PO1 (Engineering Knowledge) through PO12 (Life-long Learning), "
    "with continuous internal evaluation rubrics tracking threshold attainment targets fixed at 75% across all cohorts."
)

pdf_nba.add_table(
    headers=["NBA SAR Criterion", "Evaluation Parameter", "Max Marks", "Attained Marks", "Compliance Status"],
    rows=[
        ["Criterion 1", "Vision, Mission & Program Educational Objectives", "50", "48", "Complied"],
        ["Criterion 2", "Program Curriculum & Teaching-Learning Processes", "100", "95", "Complied"],
        ["Criterion 3", "Course Outcomes (COs) and Program Outcomes (POs)", "175", "164", "Complied"],
        ["Criterion 4", "Students' Performance & Progression Metrics", "100", "92", "Complied"],
        ["Criterion 5", "Faculty Information & Contributions (Cadre/Ret.)", "200", "188", "Complied"],
        ["Criterion 6", "Facilities and Technical Laboratories", "80", "78", "Complied"],
        ["Criterion 7", "Continuous Improvement & Indirect Feedback", "75", "72", "Complied"],
    ],
    col_widths=[30, 80, 25, 25, 30]
)

pdf_nba.add_section(
    "2. Faculty Cadre Ratio, Retention and Industrial Consulting",
    "Faculty cadre proportion adheres to AICTE guidelines with a ratio of 1 Professor : 2 Associate Professors : 6 Assistant "
    "Professors. The faculty retention rate exceeds 91.5% over a rolling 4-year assessment window. Department faculty actively "
    "lead paid industrial consulting projects with Samsung R&D, Robert Bosch, Intel, and Synopsys."
)

pdf_nba.output(os.path.join(RAW_REGULATORY, "NBA_Criteria_Compliance_Report.pdf"))
print("  [+] Generated 'NAAC_Self_Study_Summary.pdf' and 'NBA_Criteria_Compliance_Report.pdf'")


# ==============================================================================
# 4. GENERATE 16:9 WIDESCREEN PPTX PRESENTATIONS (data/raw/presentations/)
# ==============================================================================
print("[*] Generating 16:9 Widescreen Center of Excellence & Lab Showcase PPTX decks...")

prs = Presentation()
# Set widescreen 16:9 aspect ratio (13.333 x 7.5 inches)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_slide_layout = prs.slide_layouts[6]  # Completely blank layout

# Helper functions for PPTX aesthetics
def create_solid_background(slide, color_rgb):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color_rgb
    shape.line.fill.background()
    return shape

def add_header_banner(slide, category_text, title_text):
    # Top Accent Bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(37, 99, 235)  # Royal Blue
    accent.line.fill.background()

    # Category Tag
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = RGBColor(37, 99, 235)

    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(15, 23, 42)


# --- SLIDE 1: Title & Vision Deck ---
slide1 = prs.slides.add_slide(blank_slide_layout)
create_solid_background(slide1, RGBColor(15, 23, 42))  # Dark Slate

# Title Box
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.5))
tf1 = title_box.text_frame
tf1.word_wrap = True

p_sub = tf1.paragraphs[0]
p_sub.text = "PRAGYANAI INSTITUTIONAL INTELLIGENCE HUB"
p_sub.font.size = Pt(13)
p_sub.font.bold = True
p_sub.font.color.rgb = RGBColor(56, 189, 248)  # Sky Blue

p_main = tf1.add_paragraph()
p_main.text = "Centers of Excellence, R&D Labs & Academic Infrastructure"
p_main.font.size = Pt(32)
p_main.font.bold = True
p_main.font.color.rgb = RGBColor(255, 255, 255)

p_desc = tf1.add_paragraph()
p_desc.text = "A Comprehensive Showcase of Tier-1 Research Facilities, Corporate Labs & Student Innovation Testbeds"
p_desc.font.size = Pt(14)
p_desc.font.color.rgb = RGBColor(148, 163, 184)


# --- SLIDE 2: Advanced CoE Research Facilities ---
slide2 = prs.slides.add_slide(blank_slide_layout)
create_solid_background(slide2, RGBColor(248, 250, 252))
add_header_banner(slide2, "Flagship Research Facilities", "Specialized Centers of Excellence & Industry Testbeds")

# Render 3 Feature Cards
card_specs = [
    (
        "Center of Excellence in Generative AI & HPC",
        "• High-Performance NVIDIA H100/A100 GPU compute cluster\n"
        "• Research focus: Multi-agent orchestration, RAG & LLM alignment\n"
        "• Active grants: INR 3.8 Cr sponsored by DST and industry partners\n"
        "• Hands-on immersion for CSE & AI-DS undergraduate scholars",
        Inches(0.8)
    ),
    (
        "VLSI Design & Semiconductor Testbed",
        "• Complete Cadence Virtuoso & Synopsys EDA automated toolchain\n"
        "• FPGA rapid-prototyping suites (Xilinx UltraScale+ / Intel Stratix)\n"
        "• Focus: RISC-V processor architecture and RTL verification\n"
        "• Corporate hiring pipeline: Qualcomm, Intel, Texas Instruments",
        Inches(4.8)
    ),
    (
        "Robotics, Autonomous Systems & IoT Lab",
        "• Autonomous indoor drone flight arena and LiDAR sensor suites\n"
        "• Industrial 6-axis robotic arms and digital twin simulation\n"
        "• Focus: Cyber-physical systems and ROS2 robot automation\n"
        "• Supported by Bosch and Schneider Electric CoE partnerships",
        Inches(8.8)
    )
]

for card_title, card_body, left_pos in card_specs:
    card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8), Inches(3.7), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(226, 232, 240)
    card.line.width = Pt(1.5)

    tb = slide2.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = card_title
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(29, 78, 216)
    
    p2 = tf.add_paragraph()
    p2.text = card_body
    p2.font.size = Pt(9.5)
    p2.font.color.rgb = RGBColor(71, 85, 105)


# --- SLIDE 3: Placement & Career Acceleration ---
slide3 = prs.slides.add_slide(blank_slide_layout)
create_solid_background(slide3, RGBColor(248, 250, 252))
add_header_banner(slide3, "Industry Outcomes & ROI", "Pre-Placement Acceleration, Student Bootcamps & Placement Trajectory")

metrics = [
    ("62.0 LPA", "Highest Placement Offer", Inches(0.8), RGBColor(16, 185, 129)),
    ("11.50 LPA", "Computing Median CTC", Inches(3.8), RGBColor(37, 99, 235)),
    ("88.4%", "Overall Placement Rate", Inches(6.8), RGBColor(139, 92, 246)),
    ("250+ Firms", "Annual Corporate Recruiters", Inches(9.8), RGBColor(245, 158, 11))
]

for value, label, left_pos, color in metrics:
    m_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8), Inches(2.7), Inches(1.5))
    m_box.fill.solid()
    m_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    m_box.line.color.rgb = RGBColor(226, 232, 240)

    tb = slide3.shapes.add_textbox(left_pos, Inches(1.9), Inches(2.7), Inches(1.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    pv = tf.paragraphs[0]
    pv.text = value
    pv.alignment = PP_ALIGN.CENTER
    pv.font.size = Pt(20)
    pv.font.bold = True
    pv.font.color.rgb = color
    
    pl = tf.add_paragraph()
    pl.text = label
    pl.alignment = PP_ALIGN.CENTER
    pl.font.size = Pt(9)
    pl.font.color.rgb = RGBColor(100, 116, 139)

# Bottom Detail Card
b_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.6), Inches(11.7), Inches(3.0))
b_card.fill.solid()
b_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
b_card.line.color.rgb = RGBColor(226, 232, 240)

b_tb = slide3.shapes.add_textbox(Inches(1.1), Inches(3.8), Inches(11.1), Inches(2.6))
b_tf = b_tb.text_frame
b_tf.word_wrap = True

bp1 = b_tf.paragraphs[0]
bp1.text = "Holistic Career Acceleration Strategy"
bp1.font.size = Pt(14)
bp1.font.bold = True
bp1.font.color.rgb = RGBColor(15, 23, 42)

bp2 = b_tf.add_paragraph()
bp2.text = (
    "• Structured 6-Semester Skill Ladder: Foundation C++/Python (Sem 1-2) -> Full-Stack & System Design (Sem 3-4) -> "
    "Domain Specialization in AI/VLSI & Mock Interviews (Sem 5-6).\n"
    "• Corporate Mentorship & Live Capstones: Industry-sponsored final year projects in partnership with Samsung R&D, "
    "Microsoft, and Intel.\n"
    "• Competitive Coding & Hackathon Culture: Active student chapters (ACM, IEEE, GDG) with institutional travel grants "
    "for national and international hackathons."
)
bp2.font.size = Pt(10)
bp2.font.color.rgb = RGBColor(71, 85, 105)

pptx_output_path = os.path.join(RAW_PRESENTATIONS, "COE_and_Department_Infrastructure.pptx")
prs.save(pptx_output_path)
print(f"  [+] Generated 'COE_and_Department_Infrastructure.pptx'")

print("\n==========================================================================")
print("[✔] SUCCESS: All raw documents synthesized in 'data/raw/':")
print(f"    - {RAW_BROCHURES}/Admission_Flyer_2026.pdf")
print(f"    - {RAW_BROCHURES}/Placement_ROI_Report_2026.pdf")
print(f"    - {RAW_REGULATORY}/NAAC_Self_Study_Summary.pdf")
print(f"    - {RAW_REGULATORY}/NBA_Criteria_Compliance_Report.pdf")
print(f"    - {RAW_PRESENTATIONS}/COE_and_Department_Infrastructure.pptx")
print("==========================================================================")


if __name__ == "__main__":
    pass
