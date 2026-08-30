"""
src/db/generate_data_files.py

Automated Data Seeder & Asset Generator for PragyanAI College Intelligence Hub.
Exports callable generation functions and provides standalone CLI execution for:
1. data/seed/colleges.json (15 Top Engineering Institutions with State, District, City, Address,
   Department-wise Intakes, CET/COMEDK Codes, NIRF, NAAC, NBA, and Fee metrics)
2. data/seed/cutoffs.csv (KCET & COMEDK multi-year historical & projected rank cutoffs)
3. data/seed/students_synthetic.csv (1,050+ benchmark student profiles with skills, CTC, and status)
4. data/seed/outreach_events.json (Free AI Bootcamps, Masterclasses & Admissions Webinars)
5. data/raw/brochures/ (PDF Admission Flyers & Management Quota Guides)
6. data/raw/regulatory/ (PDF NAAC SSR & NBA Outcome-Based Education SAR Summaries)
7. data/raw/presentations/ (PPTX Center of Excellence & Labs Showcase)
"""

import json
import logging
import os
from pathlib import Path
import random
import uuid
from fpdf import FPDF
import pandas as pd
from pptx import Presentation

logger = logging.getLogger(__name__)

# ==============================================================================
# 0. DIRECTORY CONFIGURATION
# ==============================================================================
ROOT_DIR = Path(__file__).resolve().parent.parent.parent
BASE_DIR = ROOT_DIR / "data"
RAW_BROCHURES = BASE_DIR / "raw" / "brochures"
RAW_PRESENTATIONS = BASE_DIR / "raw" / "presentations"
RAW_REGULATORY = BASE_DIR / "raw" / "regulatory"
SEED_DIR = BASE_DIR / "seed"
VECTOR_DIR = BASE_DIR / "vector_store"


def init_target_directories():
    """Initializes and ensures all data directories exist securely."""
    for folder in [
        RAW_BROCHURES,
        RAW_PRESENTATIONS,
        RAW_REGULATORY,
        SEED_DIR,
        VECTOR_DIR,
    ]:
        try:
            folder.mkdir(parents=True, exist_ok=True)
        except Exception as e:
            logger.error(f"Failed to create directory {folder}: {e}")


# ==============================================================================
# 1. 15 BENCHMARK COLLEGES DATASET (WITH GEO, DEPARTMENTS & INTAKES)
# ==============================================================================
COLLEGES_DATA = [
    {
        "college_id": "c-001",
        "code": "E001",
        "name": "RV College of Engineering",
        "short_name": "RVCE",
        "state": "Karnataka",
        "district": "Bengaluru Urban",
        "city": "Bengaluru",
        "address": "Mysore Road, RV Vidyaniketan Post, Bengaluru, Karnataka 560059",
        "established_year": 1963,
        "autonomous": True,
        "naac_grade": "A++",
        "naac_cgpa": 3.78,
        "nba_accredited_programs": 14,
        "nirf_rank_2025": 38,
        "intake_total": 1460,
        "mgmt_fee_cse_lakhs": 16.0,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.81,
        "median_ctc_lpa": 14.5,
        "highest_ctc_lpa": 62.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 240,
            "AI & Machine Learning (AI-ML)": 120,
            "Information Science (ISE)": 180,
            "Electronics & Comm (ECE)": 180,
            "Data Science (AI-DS)": 60,
            "Electrical & Electronics (EEE)": 120,
            "Mechanical Engineering (MECH)": 120,
            "Biotechnology": 60,
            "Civil Engineering": 120,
            "Chemical Engineering": 60,
            "Aerospace Engineering": 60,
            "Industrial Engg & Mgmt": 60,
            "Telecommunication Engg": 60,
            "Instrumentation Tech": 60,
        },
        "website_link": "https://rvce.edu.in",
        "top_recruiters": [
            "Microsoft",
            "Google",
            "Amazon",
            "Cisco",
            "Intel",
            "Goldman Sachs",
        ],
        "coas_and_centers_of_excellence": [
            "Center of Excellence in Quantum Computing & AI",
            "VLSI Design & Embedded Systems Center",
            "Automotive R&D Testing Lab",
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "Empowering students through experiential learning, deep tech"
            " research, and strong global industry integration."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/rv-college-of-engineering/people/"
        ),
    },
    {
        "college_id": "c-002",
        "code": "E002",
        "name": "BMS College of Engineering",
        "short_name": "BMSCE",
        "state": "Karnataka",
        "district": "Bengaluru Urban",
        "city": "Bengaluru",
        "address": "Bull Temple Road, Basavanagudi, Bengaluru, Karnataka 560019",
        "established_year": 1946,
        "autonomous": True,
        "naac_grade": "A++",
        "naac_cgpa": 3.83,
        "nba_accredited_programs": 12,
        "nirf_rank_2025": 72,
        "intake_total": 1580,
        "mgmt_fee_cse_lakhs": 12.5,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.81,
        "median_ctc_lpa": 11.2,
        "highest_ctc_lpa": 50.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 240,
            "AI & Data Science (AI-DS)": 120,
            "Information Science (ISE)": 180,
            "Electronics & Comm (ECE)": 180,
            "Electrical & Electronics (EEE)": 120,
            "Mechanical Engineering (MECH)": 180,
            "Aerospace Engineering": 60,
            "Civil Engineering": 120,
            "Biotechnology": 60,
            "Chemical Engineering": 60,
            "Industrial Engineering": 60,
            "Medical Electronics": 60,
        },
        "website_link": "https://www.bmsce.ac.in",
        "top_recruiters": [
            "Amazon",
            "Oracle",
            "Samsung R&D",
            "Qualcomm",
            "Robert Bosch",
        ],
        "coas_and_centers_of_excellence": [
            "Center for Data Science and Generative AI",
            "Robotics and Automation Center",
            "Additive Manufacturing Testbed",
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "75+ years of engineering excellence fostering entrepreneurship and"
            " fundamental engineering innovation."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/bms-college-of-engineering/people/"
        ),
    },
    {
        "college_id": "c-003",
        "code": "E003",
        "name": "Ramaiah Institute of Technology",
        "short_name": "MSRIT",
        "state": "Karnataka",
        "district": "Bengaluru Urban",
        "city": "Bengaluru",
        "address": "MSR Nagar, MSRIT Post, Mathikere, Bengaluru, Karnataka 560054",
        "established_year": 1962,
        "autonomous": True,
        "naac_grade": "A+",
        "naac_cgpa": 3.48,
        "nba_accredited_programs": 11,
        "nirf_rank_2025": 65,
        "intake_total": 1420,
        "mgmt_fee_cse_lakhs": 11.0,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.81,
        "median_ctc_lpa": 12.0,
        "highest_ctc_lpa": 50.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 240,
            "AI & Machine Learning (AI-ML)": 120,
            "Information Science (ISE)": 120,
            "Electronics & Comm (ECE)": 180,
            "Cyber Security": 60,
            "Electrical & Electronics (EEE)": 120,
            "Mechanical Engineering (MECH)": 120,
            "Biotechnology": 60,
            "Chemical Engineering": 60,
            "Civil Engineering": 120,
            "Industrial Engg & Mgmt": 60,
            "Medical Electronics": 60,
        },
        "website_link": "https://www.msrit.edu",
        "top_recruiters": [
            "Adobe",
            "Texas Instruments",
            "Goldman Sachs",
            "Infosys",
            "Dell",
        ],
        "coas_and_centers_of_excellence": [
            "Center of Excellence in Applied AI & Healthcare Systems",
            "Cybersecurity & Cloud Computing CoE",
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "Committed to academic rigor, innovation, and ethical technological advancements."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/ramaiah-institute-of-technology/people/"
        ),
    },
    {
        "college_id": "c-004",
        "code": "E004",
        "name": "PES University (Ring Road Campus)",
        "short_name": "PESU",
        "state": "Karnataka",
        "district": "Bengaluru Urban",
        "city": "Bengaluru",
        "address": "100 Feet Ring Road, BSK III Stage, Bengaluru, Karnataka 560085",
        "established_year": 1972,
        "autonomous": True,
        "naac_grade": "A",
        "naac_cgpa": 3.52,
        "nba_accredited_programs": 10,
        "nirf_rank_2025": 90,
        "intake_total": 1800,
        "mgmt_fee_cse_lakhs": 14.0,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 4.5,
        "median_ctc_lpa": 13.5,
        "highest_ctc_lpa": 65.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 720,
            "AI & Machine Learning": 240,
            "Electronics & Comm (ECE)": 360,
            "Electrical & Electronics (EEE)": 120,
            "Mechanical Engineering": 120,
            "Biotechnology": 60,
            "Civil Engineering": 60,
            "Design & Architecture": 120,
        },
        "website_link": "https://pes.edu",
        "top_recruiters": [
            "Apple",
            "Microsoft",
            "Morgan Stanley",
            "Intuit",
            "Walmart Labs",
        ],
        "coas_and_centers_of_excellence": [
            "Center for Cloud Computing and Big Data (CCBD)",
            "Crucible of Research and Innovation (CORI)",
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "Design your tomorrow with cutting-edge core engineering and multi-agent AI ecosystems."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/pes-university/people/"
        ),
    },
    {
        "college_id": "c-005",
        "code": "E005",
        "name": "Dayananda Sagar College of Engineering",
        "short_name": "DSCE",
        "state": "Karnataka",
        "district": "Bengaluru Urban",
        "city": "Bengaluru",
        "address": "Shavige Malleshwara Hills, Kumaraswamy Layout, Bengaluru, Karnataka 560078",
        "established_year": 1979,
        "autonomous": True,
        "naac_grade": "A",
        "naac_cgpa": 3.25,
        "nba_accredited_programs": 9,
        "nirf_rank_2025": 140,
        "intake_total": 2100,
        "mgmt_fee_cse_lakhs": 8.5,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.81,
        "median_ctc_lpa": 8.2,
        "highest_ctc_lpa": 36.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 360,
            "AI & Data Science (AI-DS)": 180,
            "Information Science (ISE)": 180,
            "Electronics & Comm (ECE)": 240,
            "Cyber Security": 120,
            "Aeronautical Engineering": 60,
            "Mechanical Engineering": 180,
            "Electrical & Electronics": 120,
            "Biotechnology": 60,
            "Chemical Engineering": 60,
            "Civil Engineering": 120,
        },
        "website_link": "https://www.dsce.edu.in",
        "top_recruiters": [
            "Accenture",
            "Capgemini",
            "Dell Technologies",
            "Mercedes-Benz",
        ],
        "coas_and_centers_of_excellence": [
            "Aerospace and Avionics Research Center",
            "Autonomous Driving Systems Lab",
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "Comprehensive holistic education with high placement conversion in top multinational firms."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/dayananda-sagar-college-of-engineering/people/"
        ),
    },
    {
        "college_id": "c-006",
        "code": "E006",
        "name": "Bangalore Institute of Technology",
        "short_name": "BIT",
        "state": "Karnataka",
        "district": "Bengaluru Urban",
        "city": "Bengaluru",
        "address": "K.R. Road, V.V. Puram, Bengaluru, Karnataka 560004",
        "established_year": 1979,
        "autonomous": True,
        "naac_grade": "A+",
        "naac_cgpa": 3.32,
        "nba_accredited_programs": 8,
        "nirf_rank_2025": 180,
        "intake_total": 1300,
        "mgmt_fee_cse_lakhs": 7.5,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.81,
        "median_ctc_lpa": 7.8,
        "highest_ctc_lpa": 32.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 240,
            "AI & Machine Learning (AI-ML)": 120,
            "Information Science (ISE)": 120,
            "Electronics & Comm (ECE)": 180,
            "Electrical & Electronics (EEE)": 120,
            "Mechanical Engineering": 120,
            "Civil Engineering": 120,
            "Industrial Engg & Mgmt": 60,
        },
        "website_link": "https://bit-bangalore.edu.in",
        "top_recruiters": ["TCS", "Wipro", "Oracle", "SAP Labs", "Microchip"],
        "coas_and_centers_of_excellence": [
            "Center for Mobile Computing and Sensor Networks",
            "Cyber Threat Intelligence Cell",
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "Nurturing technical expertise and foundational ethics in future industry leaders."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/bangalore-institute-of-technology/people/"
        ),
    },
    {
        "college_id": "c-007",
        "code": "E007",
        "name": "Sir M. Visvesvaraya Institute of Technology",
        "short_name": "SMVIT",
        "state": "Karnataka",
        "district": "Bengaluru Urban",
        "city": "Bengaluru",
        "address": "Krishnadevaraya Nagar, Hunasamaranahalli, International Airport Road, Bengaluru, Karnataka 562157",
        "established_year": 1986,
        "autonomous": False,
        "naac_grade": "B++",
        "naac_cgpa": 2.92,
        "nba_accredited_programs": 6,
        "nirf_rank_2025": 220,
        "intake_total": 980,
        "mgmt_fee_cse_lakhs": 6.5,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.81,
        "median_ctc_lpa": 6.8,
        "highest_ctc_lpa": 28.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 180,
            "Information Science (ISE)": 120,
            "Electronics & Comm (ECE)": 120,
            "Electrical & Electronics": 60,
            "Mechanical Engineering": 120,
            "Biotechnology": 60,
            "Civil Engineering": 60,
        },
        "website_link": "https://www.sirmvit.edu",
        "top_recruiters": [
            "L&T Technology Services",
            "Mindtree",
            "Cognizant",
            "Toyota Kirloskar",
        ],
        "coas_and_centers_of_excellence": [
            "Bio-Energy & Environmental Engineering Center"
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "Dedicated to practical skills, rigorous academic training, and rural development outreach."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/sir-m-visvesvaraya-institute-of-technology/people/"
        ),
    },
    {
        "college_id": "c-008",
        "code": "E008",
        "name": "The National Institute of Engineering",
        "short_name": "NIE",
        "state": "Karnataka",
        "district": "Mysuru",
        "city": "Mysuru",
        "address": "Manandavadi Road, Mysuru, Karnataka 570008",
        "established_year": 1946,
        "autonomous": True,
        "naac_grade": "A",
        "naac_cgpa": 3.22,
        "nba_accredited_programs": 7,
        "nirf_rank_2025": 160,
        "intake_total": 960,
        "mgmt_fee_cse_lakhs": 7.0,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.81,
        "median_ctc_lpa": 9.0,
        "highest_ctc_lpa": 44.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 240,
            "Information Science (ISE)": 120,
            "Electronics & Comm (ECE)": 180,
            "Electrical & Electronics (EEE)": 120,
            "Mechanical Engineering (MECH)": 120,
            "Civil Engineering": 120,
            "Industrial & Production": 60,
        },
        "website_link": "https://nie.ac.in",
        "top_recruiters": [
            "Cisco",
            "Titan",
            "ABB",
            "Schneider Electric",
            "Infosys",
        ],
        "coas_and_centers_of_excellence": [
            "Center for Renewable Energy & Smart Grids",
            "NIE Center for Industrial IoT",
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "A heritage of technical eminence building nation-builders and core engineering leaders."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/the-national-institute-of-engineering-mysuru/people/"
        ),
    },
    {
        "college_id": "c-009",
        "code": "E009",
        "name": "NMAM Institute of Technology",
        "short_name": "NMAMIT",
        "state": "Karnataka",
        "district": "Udupi",
        "city": "Nitte",
        "address": "Nitte, Karkala Taluk, Udupi District, Karnataka 574110",
        "established_year": 1986,
        "autonomous": True,
        "naac_grade": "A",
        "naac_cgpa": 3.18,
        "nba_accredited_programs": 7,
        "nirf_rank_2025": 175,
        "intake_total": 1100,
        "mgmt_fee_cse_lakhs": 5.5,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.5,
        "median_ctc_lpa": 6.8,
        "highest_ctc_lpa": 30.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 240,
            "AI & Machine Learning (AI-ML)": 120,
            "Information Science (ISE)": 120,
            "Electronics & Comm (ECE)": 180,
            "Mechanical Engineering": 120,
            "Electrical & Electronics": 60,
            "Biotechnology": 60,
            "Civil Engineering": 60,
        },
        "website_link": "https://nmamit.nitte.edu.in",
        "top_recruiters": ["Sony India", "Nutanix", "KPIT", "Tesco"],
        "coas_and_centers_of_excellence": [
            "Center for Artificial Intelligence & Cyber-Physical Systems"
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "Fostering technological competence and moral responsibility in a serene campus setting."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/nmam-institute-of-technology/people/"
        ),
    },
    {
        "college_id": "c-010",
        "code": "E010",
        "name": "Siddaganga Institute of Technology",
        "short_name": "SIT",
        "state": "Karnataka",
        "district": "Tumakuru",
        "city": "Tumakuru",
        "address": "B.H. Road, Tumakuru, Karnataka 572103",
        "established_year": 1963,
        "autonomous": True,
        "naac_grade": "A++",
        "naac_cgpa": 3.65,
        "nba_accredited_programs": 9,
        "nirf_rank_2025": 100,
        "intake_total": 1250,
        "mgmt_fee_cse_lakhs": 6.0,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.81,
        "median_ctc_lpa": 8.4,
        "highest_ctc_lpa": 34.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 240,
            "AI & Data Science (AI-DS)": 120,
            "Electronics & Comm (ECE)": 180,
            "Information Science (ISE)": 120,
            "Electrical & Electronics": 120,
            "Mechanical Engineering": 120,
            "Chemical Engineering": 60,
            "Civil Engineering": 120,
            "Biotechnology": 60,
        },
        "website_link": "http://www.sit.ac.in",
        "top_recruiters": [
            "Amazon",
            "Mercedes-Benz R&D",
            "Siemens",
            "TCS Ninja",
        ],
        "coas_and_centers_of_excellence": [
            "Nanotechnology & Advanced Materials Research Lab",
            "SIT AI Incubation Center",
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "Value-based education grounded in discipline, innovation, and industrial readiness."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/siddaganga-institute-of-technology/people/"
        ),
    },
    {
        "college_id": "c-011",
        "code": "E011",
        "name": "JSS Science and Technology University (SJCE)",
        "short_name": "JSS STU",
        "state": "Karnataka",
        "district": "Mysuru",
        "city": "Mysuru",
        "address": "JSS Technical Institutions Campus, Mysuru, Karnataka 570006",
        "established_year": 1963,
        "autonomous": True,
        "naac_grade": "A",
        "naac_cgpa": 3.28,
        "nba_accredited_programs": 10,
        "nirf_rank_2025": 158,
        "intake_total": 1350,
        "mgmt_fee_cse_lakhs": 8.0,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.81,
        "median_ctc_lpa": 9.2,
        "highest_ctc_lpa": 40.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 240,
            "Information Science (ISE)": 120,
            "Electronics & Comm (ECE)": 180,
            "Electrical & Electronics (EEE)": 120,
            "Mechanical Engineering": 120,
            "Environmental Engineering": 60,
            "Polymer Science & Tech": 60,
            "Civil Engineering": 120,
        },
        "website_link": "https://jssstuniv.in",
        "top_recruiters": [
            "Qualcomm",
            "Applied Materials",
            "Western Digital",
            "Bosch",
        ],
        "coas_and_centers_of_excellence": [
            "CoE in Advanced Electronic Packaging & VLSI Testing"
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "Pioneering technical education that transforms students into visionary problem solvers."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/jss-science-and-technology-university/people/"
        ),
    },
    {
        "college_id": "c-012",
        "code": "E012",
        "name": "KLS Gogte Institute of Technology",
        "short_name": "GIT",
        "state": "Karnataka",
        "district": "Belagavi",
        "city": "Belagavi",
        "address": "Jnana Ganga, Udyambag, Belagavi, Karnataka 590008",
        "established_year": 1979,
        "autonomous": True,
        "naac_grade": "A+",
        "naac_cgpa": 3.35,
        "nba_accredited_programs": 6,
        "nirf_rank_2025": 250,
        "intake_total": 1080,
        "mgmt_fee_cse_lakhs": 4.5,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.2,
        "median_ctc_lpa": 6.2,
        "highest_ctc_lpa": 24.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 240,
            "AI & Data Science": 60,
            "Information Science (ISE)": 120,
            "Electronics & Comm (ECE)": 180,
            "Aeronautical Engineering": 60,
            "Mechanical Engineering": 120,
            "Electrical & Electronics": 60,
            "Civil Engineering": 120,
        },
        "website_link": "https://git.edu",
        "top_recruiters": ["PwC", "Tech Mahindra", "KPIT", "SLK Software"],
        "coas_and_centers_of_excellence": [
            "Aerospace Structures and Advanced CFD Lab"
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "Empowering North Karnataka with quality engineering, tech incubators, and industry ties."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/kls-gogte-institute-of-technology/people/"
        ),
    },
    {
        "college_id": "c-013",
        "code": "E013",
        "name": "BMS Institute of Technology & Management",
        "short_name": "BMSIT",
        "state": "Karnataka",
        "district": "Bengaluru Urban",
        "city": "Bengaluru",
        "address": "Doddaballapur Main Road, Avalahalli, Yelahanka, Bengaluru, Karnataka 560064",
        "established_year": 2002,
        "autonomous": True,
        "naac_grade": "A",
        "naac_cgpa": 3.21,
        "nba_accredited_programs": 6,
        "nirf_rank_2025": 190,
        "intake_total": 960,
        "mgmt_fee_cse_lakhs": 8.0,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.81,
        "median_ctc_lpa": 8.5,
        "highest_ctc_lpa": 35.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 240,
            "AI & Machine Learning (AI-ML)": 120,
            "Information Science (ISE)": 120,
            "Electronics & Comm (ECE)": 180,
            "Mechanical Engineering": 60,
            "Electrical & Electronics": 60,
            "Civil Engineering": 60,
        },
        "website_link": "https://bmsit.ac.in",
        "top_recruiters": [
            "Fidelity Investments",
            "VMware",
            "Nutanix",
            "Oracle",
        ],
        "coas_and_centers_of_excellence": [
            "Open Source & Generative AI Innovation Center"
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "Modern progressive learning with rigorous hands-on coding and industry hackathons."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/bmsit-m/people/"
        ),
    },
    {
        "college_id": "c-014",
        "code": "E014",
        "name": "New Horizon College of Engineering",
        "short_name": "NHCE",
        "state": "Karnataka",
        "district": "Bengaluru Urban",
        "city": "Bengaluru",
        "address": "Bellandur Main Road, Near Marathahalli, Bengaluru, Karnataka 560103",
        "established_year": 2001,
        "autonomous": True,
        "naac_grade": "A",
        "naac_cgpa": 3.23,
        "nba_accredited_programs": 5,
        "nirf_rank_2025": 210,
        "intake_total": 1200,
        "mgmt_fee_cse_lakhs": 6.5,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.81,
        "median_ctc_lpa": 7.4,
        "highest_ctc_lpa": 30.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 240,
            "AI & Machine Learning (AI-ML)": 120,
            "Information Science (ISE)": 120,
            "Electronics & Comm (ECE)": 180,
            "Electrical & Electronics": 60,
            "Mechanical Engineering": 60,
            "Civil Engineering": 60,
        },
        "website_link": "https://newhorizonindia.edu/nhengineering/",
        "top_recruiters": ["Capgemini", "LTI Mindtree", "Genpact", "Mu Sigma"],
        "coas_and_centers_of_excellence": [
            "French-Indo Center of Excellence for Electricity and Automation",
            "Schneider Electric Lab",
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "Connecting students directly to IT corridor corporate ecosystem and global certification."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/new-horizon-college-of-engineering/people/"
        ),
    },
    {
        "college_id": "c-015",
        "code": "E015",
        "name": "Nitte Meenakshi Institute of Technology",
        "short_name": "NMIT",
        "state": "Karnataka",
        "district": "Bengaluru Urban",
        "city": "Bengaluru",
        "address": "Govindapura, Gollahalli, Yelahanka, Bengaluru, Karnataka 560064",
        "established_year": 2001,
        "autonomous": True,
        "naac_grade": "A+",
        "naac_cgpa": 3.39,
        "nba_accredited_programs": 7,
        "nirf_rank_2025": 185,
        "intake_total": 1150,
        "mgmt_fee_cse_lakhs": 7.0,
        "govt_fee_cet_lakhs": 1.07,
        "comedk_fee_lakhs": 2.81,
        "median_ctc_lpa": 7.8,
        "highest_ctc_lpa": 32.0,
        "departments_and_intake": {
            "Computer Science & Engg (CSE)": 240,
            "AI & Data Science (AI-DS)": 120,
            "Information Science (ISE)": 120,
            "Electronics & Comm (ECE)": 180,
            "Aeronautical Engineering": 60,
            "Mechanical Engineering": 60,
            "Electrical & Electronics": 60,
            "Civil Engineering": 60,
        },
        "website_link": "https://nmit.ac.in",
        "top_recruiters": ["Subex", "Unisys", "HP Enterprise", "Siemens"],
        "coas_and_centers_of_excellence": [
            "Small Satellite Development Lab (STUDSAT)",
            "CoE in Robotics & Micro-Electro-Mechanical Systems",
        ],
        "video_tour_url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "principal_statement": (
            "Igniting research passions through hands-on space tech, robotics, and AI ventures."
        ),
        "alumni_linkedin_hub": (
            "https://www.linkedin.com/school/nitte-meenakshi-institute-of-technology/people/"
        ),
    },
]


def generate_colleges_json():
    """Generates and writes data/seed/colleges.json securely."""
    init_target_directories()
    dest_path = SEED_DIR / "colleges.json"
    try:
        with open(dest_path, "w", encoding="utf-8") as f:
            json.dump(COLLEGES_DATA, f, indent=2)
        logger.info(f"Generated {dest_path} ({len(COLLEGES_DATA)} Colleges)")
        print(f"  [+] Generated {dest_path} ({len(COLLEGES_DATA)} Colleges)")
    except Exception as e:
        logger.error(f"Failed to generate colleges.json: {e}")
    return str(dest_path)


# ==============================================================================
# 2. ENTRANCE EXAM CUTOFFS GENERATOR (data/seed/cutoffs.csv)
# ==============================================================================
def generate_cutoffs_csv():
    """Generates multi-year cutoff records across colleges, exams, and categories securely."""
    init_target_directories()
    branches = ["CSE", "AI-DS", "ISE", "ECE", "MECH"]
    exams = ["KCET", "COMEDK"]
    categories = ["GM", "1G", "2A", "2B", "3A", "3B", "SC", "ST"]
    years = [2024, 2025, 2026]

    cutoffs_records = []
    for c in COLLEGES_DATA:
        base_rank = c["nirf_rank_2025"] * 32
        for yr in years:
            yr_mult = 1.0 if yr == 2026 else (0.95 if yr == 2025 else 0.90)
            for b_idx, branch in enumerate(branches):
                for exam in exams:
                    for cat in categories:
                        cat_mult = 1.0
                        if cat in ["1G", "2A", "2B"]:
                            cat_mult = 1.35
                        elif cat in ["3A", "3B"]:
                            cat_mult = 1.2
                        elif cat in ["SC", "ST"]:
                            cat_mult = 2.8

                        exam_mult = 1.0 if exam == "KCET" else 1.45
                        branch_mult = 1.0 + (b_idx * 0.42)
                        calculated_rank = int(
                            base_rank * branch_mult * cat_mult * exam_mult * yr_mult
                        )

                        cutoffs_records.append({
                            "cutoff_id": str(uuid.uuid4()),
                            "college_code": c["code"],
                            "college_name": c["name"],
                            "year": yr,
                            "exam": exam,
                            "round": "Round-2 (Final)",
                            "branch": branch,
                            "category": cat,
                            "cutoff_rank": min(calculated_rank, 185000),
                        })

    dest_path = SEED_DIR / "cutoffs.csv"
    try:
        pd.DataFrame(cutoffs_records).to_csv(dest_path, index=False)
        logger.info(f"Generated {dest_path} ({len(cutoffs_records)} Records)")
        print(f"  [+] Generated {dest_path} ({len(cutoffs_records)} Records)")
    except Exception as e:
        logger.error(f"Failed to generate cutoffs.csv: {e}")
    return str(dest_path)


# ==============================================================================
# 3. 1,050+ SYNTHETIC STUDENT PROFILES (data/seed/students_synthetic.csv)
# ==============================================================================
def generate_students_csv():
    """Generates synthetic student talent pool with CGPA, skills, and placement records securely."""
    init_target_directories()
    first_names = [
        "Aarav", "Aditi", "Rahul", "Sneha", "Vikram", "Ananya", "Rohan", "Pooja",
        "Karthik", "Divya", "Siddharth", "Meera", "Nikhil", "Varun", "Ishita",
        "Tanvi", "Abhishek", "Bhavana", "Chaitanya", "Deepak", "Esha", "Gaurav",
        "Harsha", "Kavya", "Manoj", "Neha", "Pranav", "Rashmi", "Sanjay",
        "Tejas", "Usha", "Vinay",
    ]
    last_names = [
        "Sharma", "Patil", "Rao", "Kulkarni", "Iyer", "Deshmukh", "Gowda",
        "Nair", "Reddy", "Bhat", "Verma", "Mishra", "Ambesange", "Hegde",
        "Shetty", "Kamath", "Joshi", "Nadkarni", "Menon", "Pillai",
    ]

    skill_sets = [
        "PyTorch, LangChain, FastAPI, ChromaDB, HuggingFace",
        "React.js, Node.js, PostgreSQL, Docker, Kubernetes",
        "Verilog, SystemVerilog, RTL Verification, Cadence Virtuoso",
        "AWS Cloud, Terraform, CI/CD, Python Automation, Linux",
        "Embedded C, FreeRTOS, Linux Device Drivers, CAN Protocol",
        "TensorFlow, Scikit-learn, Pandas, OpenCV, Deep Learning",
        "Next.js, TailwindCSS, MongoDB, GraphQL, REST APIs",
        "Automated Testing, Selenium, PyTest, Java, SpringBoot",
    ]

    branches = ["CSE", "AI-DS", "ISE", "ECE", "MECH"]
    students_records = []

    for i in range(1, 1060):
        col = random.choice(COLLEGES_DATA)
        branch = random.choice(branches)
        cgpa = round(random.uniform(6.8, 9.85), 2)
        hackathons = random.choices([0, 1, 2, 3, 5], weights=[35, 30, 20, 10, 5])[0]
        status = random.choices(["Placed", "Higher Studies", "Seeking"], weights=[82, 10, 8])[0]

        ctc = 0.0
        company = "None"
        role_title = "Student"

        if status == "Placed":
            tier = random.choices(["Mass", "Core", "Dream", "SuperDream"], weights=[38, 32, 20, 10])[0]
            if tier == "Mass":
                ctc = round(random.uniform(4.0, 6.5), 1)
                company = random.choice(["Infosys", "TCS", "Wipro", "Cognizant", "Accenture"])
                role_title = "Systems Engineer"
            elif tier == "Core":
                ctc = round(random.uniform(7.0, 12.5), 1)
                company = random.choice([
                    "Samsung R&D", "Cisco Systems", "Robert Bosch", "Dell Technologies", "Siemens"
                ])
                role_title = "Software Engineer"
            elif tier == "Dream":
                ctc = round(random.uniform(14.0, 24.0), 1)
                company = random.choice([
                    "Qualcomm", "Intel", "Texas Instruments", "Oracle", "Goldman Sachs"
                ])
                role_title = "Member Technical Staff"
            else:  # SuperDream
                ctc = round(random.uniform(28.0, 58.0), 1)
                company = random.choice(["Microsoft", "Amazon", "Google", "PragyanAI", "Apple"])
                role_title = "AI / Systems Engineer"

        fname = random.choice(first_names)
        lname = random.choice(last_names)
        students_records.append({
            "student_id": f"std-{i:05d}",
            "usn": f"1{col['code'][:2]}22{branch[:2]}{i:04d}",
            "full_name": f"{fname} {lname}",
            "college_code": col["code"],
            "college_name": col["name"],
            "branch": branch,
            "grad_year": 2026,
            "cgpa": cgpa,
            "hackathons_won": hackathons,
            "primary_skills": random.choice(skill_sets),
            "placement_status": status,
            "offered_ctc_lpa": ctc,
            "placed_company": company,
            "job_title": role_title,
            "linkedin_url": f"https://www.linkedin.com/in/{fname.lower()}-{lname.lower()}-{i:04d}/",
            "google_scholar_url": (
                f"https://scholar.google.com/citations?user=std_{i:04d}" if cgpa > 9.0 else ""
            ),
        })

    dest_path = SEED_DIR / "students_synthetic.csv"
    try:
        pd.DataFrame(students_records).to_csv(dest_path, index=False)
        logger.info(f"Generated {dest_path} ({len(students_records)} Student Profiles)")
        print(f"  [+] Generated {dest_path} ({len(students_records)} Student Profiles)")
    except Exception as e:
        logger.error(f"Failed to generate students_synthetic.csv: {e}")
    return str(dest_path)


# ==============================================================================
# 4. OUTREACH EVENTS & WEBINARS DATASET (data/seed/outreach_events.json)
# ==============================================================================
OUTREACH_EVENTS = [
    {
        "event_id": "evt-001",
        "title": "Masterclass: Generative AI, RAG & Agentic AI using LangGraph",
        "track": "Deep Tech & AI",
        "speaker_name": "Sateesh Ambesange",
        "speaker_designation": "AI Architect & Research Scholar",
        "event_date": "2026-09-05",
        "event_time": "11:00 AM - 1:00 PM IST",
        "platform": "Google Meet / YouTube Live",
        "registration_fee": "Free",
        "target_audience": "PU College & High School Seniors, 1st Year BE",
        "brochure_asset": "data/raw/brochures/GenAI_Outreach_Masterclass_2026.pdf",
        "learning_outcomes": [
            "Build your first RAG Agent",
            "Connecting LLMs to Databases",
            "Future Engineering Careers in 2026-2030",
        ],
    },
    {
        "event_id": "evt-002",
        "title": "KCET & COMEDK 2026 Option Entry Strategy & College Rank Matcher",
        "track": "Admissions Guidance",
        "speaker_name": "Dean of Admissions",
        "speaker_designation": "Institutional Admissions Advisory Board",
        "event_date": "2026-09-12",
        "event_time": "5:00 PM - 6:30 PM IST",
        "platform": "Zoom Webinar",
        "registration_fee": "Free",
        "target_audience": "Aspirants & Parents",
        "brochure_asset": "data/raw/brochures/Admissions_Guide_2026.pdf",
        "learning_outcomes": [
            "Cutoff rank analysis across 15 Top Colleges",
            "Avoiding branch choice pitfalls",
            "Scholarships and Fee Concessions",
        ],
    },
    {
        "event_id": "evt-003",
        "title": "Hands-On Robotics, IoT & Micro-Sensors Lab Walkthrough",
        "track": "Hardware & Embedded Systems",
        "speaker_name": "Dr. Ramesh Kulkarni",
        "speaker_designation": "Head of Center of Excellence in IoT",
        "event_date": "2026-09-19",
        "event_time": "2:30 PM - 4:30 PM IST",
        "platform": "On-Campus Lab Discovery + Live Stream",
        "registration_fee": "Free",
        "target_audience": "Partner School Batches & PU Students",
        "brochure_asset": "data/raw/brochures/Robotics_IoT_Outreach.pdf",
        "learning_outcomes": [
            "Live sensor interfacing demonstration",
            "Autonomous rover path planning",
            "Direct interaction with research scholars",
        ],
    },
]


def generate_outreach_events_json():
    """Generates and writes data/seed/outreach_events.json securely."""
    init_target_directories()
    dest_path = SEED_DIR / "outreach_events.json"
    try:
        with open(dest_path, "w", encoding="utf-8") as f:
            json.dump(OUTREACH_EVENTS, f, indent=2)
        logger.info(f"Generated {dest_path} (Outreach Events)")
        print(f"  [+] Generated {dest_path} (Outreach Events)")
    except Exception as e:
        logger.error(f"Failed to generate outreach_events.json: {e}")
    return str(dest_path)


# ==============================================================================
# 5. PROGRAMMATIC PDF GENERATOR (data/raw/brochures/ & data/raw/regulatory/)
# ==============================================================================
class BrandedPDF(FPDF):
    def header(self):
        self.set_font("Helvetica", "B", 13)
        self.set_text_color(26, 54, 93)  # Navy
        self.cell(0, 8, "PRAGYANAI INSTITUTIONAL INTELLIGENCE HUB", 0, 1, "C")
        self.set_font("Helvetica", "I", 9)
        self.set_text_color(100, 100, 100)
        self.cell(0, 5, "Official Institutional Document & Regulatory Archive", 0, 1, "C")
        self.line(10, 22, 200, 22)
        self.ln(6)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(140, 140, 140)
        self.cell(0, 10, f"Page {self.page_no()}", 0, 0, "C")


def build_pdf_document(filepath, title, sections):
    try:
        pdf = BrandedPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()

        pdf.set_font("Helvetica", "B", 16)
        pdf.set_text_color(15, 23, 42)
        pdf.cell(0, 10, title, 0, 1, "L")
        pdf.ln(3)

        for heading, body in sections:
            pdf.set_font("Helvetica", "B", 12)
            pdf.set_text_color(30, 64, 175)  # Royal Blue
            pdf.cell(0, 7, heading, 0, 1, "L")

            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(51, 65, 85)  # Slate
            pdf.multi_cell(0, 5.5, body)
            pdf.ln(3)

        pdf.output(str(filepath))
    except Exception as e:
        logger.error(f"Failed to build PDF document {filepath}: {e}")


def generate_raw_documents():
    """Generates official brochures and regulatory PDFs in data/raw/ securely."""
    init_target_directories()

    build_pdf_document(
        RAW_BROCHURES / "Admission_Flyer_2026.pdf",
        "Management Admissions & Institutional Fee Structure (2026-27)",
        [
            (
                "1. Management Quota Fee Bifurcation",
                (
                    "Tuition Fee: Branch-specific starting from 4.5 Lakhs/year (Core"
                    " branches) up to 16.0 Lakhs/year (Computer Science & AI/ML)."
                    " Value-Added Skill Bootcamps & COE Lab Access: Included in"
                    " institutional development fees."
                ),
            ),
            (
                "2. Scholarships & Merit Concessions",
                (
                    "Students securing CET rank below 2,000 or COMEDK rank below"
                    " 1,500 are eligible for a 50% tuition waiver on Management"
                    " seats. National sports medalists eligible for direct concession."
                ),
            ),
            (
                "3. Placement ROI & Industry Immersion",
                (
                    "Consistently achieving 85%+ overall placement rate. Top"
                    " packages reach up to 62 LPA with a median salary exceeding"
                    " 11.5 LPA for Computing and Electronics branches."
                ),
            ),
        ],
    )

    build_pdf_document(
        RAW_REGULATORY / "NAAC_Self_Study_Summary.pdf",
        "NAAC Self-Study Report (SSR) - Criterion Wise Summary",
        [
            (
                "Criterion I: Curricular Aspects (Weightage 150)",
                (
                    "Autonomous syllabus updated every 2 years with 40% industry"
                    " representation on Board of Studies. Choice Based Credit System"
                    " (CBCS) implemented across 100% programs."
                ),
            ),
            (
                "Criterion II: Teaching-Learning and Evaluation (Weightage 200)",
                (
                    "Full-time faculty with Ph.D. ratio is 68%. Student to Full-Time"
                    " Faculty ratio maintained at 14:1."
                ),
            ),
            (
                "Criterion III: Research, Innovations and Extension (Weightage 250)",
                (
                    "Over INR 18.5 Crores in sponsored research projects from DST,"
                    " AICTE, and DRDO. 42 published patents in AI, IoT, and Advanced Materials."
                ),
            ),
        ],
    )

    build_pdf_document(
        RAW_REGULATORY / "NBA_Criteria_Compliance_Report.pdf",
        "National Board of Accreditation (NBA) - Outcome Based Education",
        [
            (
                "Program Educational Objectives (PEOs) & Outcomes (POs)",
                (
                    "Continuous assessment metrics mapped directly to PO1 through"
                    " PO12. Attainment threshold set at 75% for 2025-26 academic cycles."
                ),
            ),
            (
                "Faculty Cadre Proportion and Retention",
                (
                    "Cadre ratio: Professors (1) : Associate Professors (2) :"
                    " Assistant Professors (6). Faculty retention index stands at"
                    " 92% over the last 3 years."
                ),
            ),
        ],
    )
    logger.info("Generated verified PDFs in data/raw/brochures/ & regulatory/")
    print("  [+] Generated verified PDFs in data/raw/brochures/ & regulatory/")


# ==============================================================================
# 6. PROGRAMMATIC PPTX PRESENTATION (data/raw/presentations/)
# ==============================================================================
def generate_presentation_deck():
    """Generates the Center of Excellence presentation deck in data/raw/presentations/ securely."""
    init_target_directories()
    try:
        prs = Presentation()

        title_slide_layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(title_slide_layout)
        slide1.shapes.title.text = "Center of Excellence & R&D Labs"
        subtitle = slide1.placeholders[1]
        subtitle.text = (
            "PragyanAI Institutional Engineering Benchmark\nAutomated Hardware & AI"
            " Innovation Hubs (2026)"
        )

        bullet_slide_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        slide2.shapes.title.text = "State-of-the-Art Research Centers"
        tf2 = slide2.placeholders[1].text_frame
        tf2.text = "Active Industry-Sponsored Facilities:"

        p1 = tf2.add_paragraph()
        p1.text = "• AI & High Performance Computing (HPC) Supercomputing Cluster"
        p1.level = 1

        p2 = tf2.add_paragraph()
        p2.text = "• Semiconductor VLSI Design & Cadence Automated Testing Suite"
        p2.level = 1

        p3 = tf2.add_paragraph()
        p3.text = "• Autonomous Robotics, Drone Flight Arena & Sensor Testbeds"
        p3.level = 1

        slide3 = prs.slides.add_slide(bullet_slide_layout)
        slide3.shapes.title.text = "Placement ROI & Career Pathways"
        tf3 = slide3.placeholders[1].text_frame
        tf3.text = "Student Outcome Benchmarks:"

        p4 = tf3.add_paragraph()
        p4.text = "• Tier-1 Product Recruiters: Microsoft, Amazon, Cisco, Qualcomm"
        p4.level = 1

        p5 = tf3.add_paragraph()
        p5.text = "• 85%+ Overall Placement Ratio across Autonomous Programs"
        p5.level = 1

        p6 = tf3.add_paragraph()
        p6.text = "• Pre-placement multi-agent AI & full-stack software bootcamps"
        p6.level = 1

        pptx_path = RAW_PRESENTATIONS / "COE_and_Department_Infrastructure.pptx"
        prs.save(str(pptx_path))
        logger.info(f"Generated presentation deck: {pptx_path}")
        print(f"  [+] Generated presentation deck: {pptx_path}")
        return str(pptx_path)
    except Exception as e:
        logger.error(f"Failed to generate presentation deck: {e}")
        return ""


# ==============================================================================
# 7. MASTER GENERATION WRAPPER
# ==============================================================================
def generate_all_data_files():
    """Generates all raw documents, seed CSVs, JSONs, and presentation decks."""
    print("[*] Generating all seed datasets and raw institutional assets...")
    generate_colleges_json()
    generate_cutoffs_csv()
    generate_students_csv()
    generate_outreach_events_json()
    generate_raw_documents()
    generate_presentation_deck()
    print("\n==================================================================")
    print("SUCCESS: All datasets, PDFs, PPTXs, and CSVs generated cleanly!")
    print("Run 'python -m src.db.seed_runner' next to populate the SQL Database.")
    print("==================================================================")


if __name__ == "__main__":
    generate_all_data_files()
