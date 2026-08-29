"""
src/rag_engine/ingestion.py

Multimodal document parser and chunking pipeline using PyMuPDF (fitz) and python-pptx.
Parses Admission Brochures, NAAC SSRs, NBA SARs, and CoE Slides with strict metadata tagging.
"""

import os
from pathlib import Path
from typing import Dict, List, Optional
import fitz  # PyMuPDF
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation

from src.core.config import settings


def extract_pdf_chunks(
    pdf_path: str,
    doc_category: str = "General",
    chunk_size: int = settings.CHUNK_SIZE,
    chunk_overlap: int = settings.CHUNK_OVERLAP,
) -> List[Document]:
    """Extracts structured text from a PDF file with page-level metadata and chunks it."""
    file_path = Path(pdf_path)
    if not file_path.exists():
        return []

    doc = fitz.open(str(file_path))
    raw_documents: List[Document] = []

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        text = page.get_text()
        if text and text.strip():
            raw_documents.append(
                Document(
                    page_content=text.strip(),
                    metadata={
                        "source": file_path.name,
                        "file_path": str(file_path),
                        "file_type": "PDF",
                        "doc_category": doc_category,
                        "page": page_idx + 1,
                        "total_pages": len(doc),
                    },
                )
            )

    doc.close()

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return splitter.split_documents(raw_documents)


def extract_pptx_chunks(
    pptx_path: str,
    doc_category: str = "Presentation",
    chunk_size: int = settings.CHUNK_SIZE,
    chunk_overlap: int = settings.CHUNK_OVERLAP,
) -> List[Document]:
    """Extracts slide headers and body contents from PPTX presentation files."""
    file_path = Path(pptx_path)
    if not file_path.exists():
        return []

    prs = Presentation(str(file_path))
    raw_documents: List[Document] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        slide_texts.append(line)

        if slide_texts:
            content = f"Slide {slide_idx + 1}:\n" + "\n".join(slide_texts)
            raw_documents.append(
                Document(
                    page_content=content,
                    metadata={
                        "source": file_path.name,
                        "file_path": str(file_path),
                        "file_type": "PPTX",
                        "doc_category": doc_category,
                        "slide_number": slide_idx + 1,
                        "total_slides": len(prs.slides),
                    },
                )
            )

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
    )
    return splitter.split_documents(raw_documents)


class DocumentIngestionPipeline:
    """Master pipeline managing recursive directory scanning, extraction, and chunking."""

    def __init__(self):
        self.brochures_dir = Path(settings.BROCHURES_DIR)
        self.presentations_dir = Path(settings.PRESENTATIONS_DIR)
        self.regulatory_dir = Path(settings.REGULATORY_DIR)

    def process_all(self) -> List[Document]:
        """Scans all raw directories and converts all files into searchable Document chunks."""
        all_chunks: List[Document] = []

        # 1. Ingest Brochures & Admission Guides
        if self.brochures_dir.exists():
            for pdf_file in self.brochures_dir.glob("*.pdf"):
                chunks = extract_pdf_chunks(str(pdf_file), doc_category="Admissions & Fees")
                all_chunks.extend(chunks)

        # 2. Ingest Regulatory Dossiers (NAAC / NBA)
        if self.regulatory_dir.exists():
            for pdf_file in self.regulatory_dir.glob("*.pdf"):
                category = "NBA Accreditation" if "NBA" in pdf_file.name else "NAAC SSR Audit"
                chunks = extract_pdf_chunks(str(pdf_file), doc_category=category)
                all_chunks.extend(chunks)

        # 3. Ingest CoE Presentation Decks
        if self.presentations_dir.exists():
            for pptx_file in self.presentations_dir.glob("*.pptx"):
                chunks = extract_pptx_chunks(str(pptx_file), doc_category="Centers of Excellence")
                all_chunks.extend(chunks)

        return all_chunks


def ingest_all_raw_documents() -> List[Document]:
    """Helper method executing full directory ingestion."""
    pipeline = DocumentIngestionPipeline()
    return pipeline.process_all()
